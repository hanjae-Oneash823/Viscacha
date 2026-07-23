"""
Step 0.1 — Metadata harmonization
Reads SMC xlsx and PO pptx, harmonizes donor IDs, builds unified metadata table.
"""

import re
import pandas as pd
from pathlib import Path
from pptx import Presentation

from layer0.config import (
    SMC_XLSX, PO_PPTX, UNIFIED_META_COLS,
    COND_AD, COND_CTRL, COND_ACTIVE,
    B_TO_BRAAK, ROMAN_TO_INT, CERAD_TEXT_TO_INT,
)
from layer0.utils.qc_log import QCLogger


# ---------------------------------------------------------------------------
# ID harmonization
# ---------------------------------------------------------------------------

def harmonize_id(raw: str) -> str:
    """
    SMC-027  -> SMC027   (strip hyphen)
    PO-05    -> PO05     (strip hyphen)
    P_005    -> PO05     (replace P_ prefix, 2-digit zero-pad)
    """
    s = str(raw).strip()
    if re.match(r'^SMC-\d+$', s):
        return s.replace('SMC-', 'SMC')
    if re.match(r'^PO-\d+$', s):
        return s.replace('PO-', 'PO')
    if re.match(r'^P_\d+$', s):
        digits = s.split('_')[1]
        return 'PO' + str(int(digits)).zfill(2)
    return s   # already harmonized or unrecognized


# ---------------------------------------------------------------------------
# Condition normalization
# ---------------------------------------------------------------------------

_COND_MAP = {
    'ad':             COND_AD,
    'control':        COND_CTRL,
    'active control': COND_ACTIVE,
    'nc':             COND_CTRL,   # PO "Normal Control"
}

def _norm_condition(raw: str) -> str:
    return _COND_MAP.get(str(raw).strip().lower(), str(raw).strip())


# ---------------------------------------------------------------------------
# Helpers for ABC score parsing
# ---------------------------------------------------------------------------

def _b_to_braak(b_val) -> float:
    """Convert xlsx B column value (e.g. 'B2', 2.0) to numeric Braak midpoint."""
    if pd.isna(b_val):
        return float('nan')
    s = str(b_val).strip().upper()
    if not s.startswith('B'):
        s = 'B' + s.lstrip('0') or 'B0'
    return float(B_TO_BRAAK.get(s, float('nan')))


def _abc_to_int(val) -> float:
    """Convert A1 / B2 / C3 style to integer tier."""
    if pd.isna(val):
        return float('nan')
    s = str(val).strip()
    m = re.search(r'\d', s)
    return float(m.group()) if m else float('nan')


def _roman_to_int(val) -> float:
    if pd.isna(val):
        return float('nan')
    s = str(val).strip()
    return float(ROMAN_TO_INT.get(s, float('nan')))


def _cerad_text_to_int(val) -> float:
    if pd.isna(val):
        return float('nan')
    return float(CERAD_TEXT_TO_INT.get(str(val).strip().lower(), float('nan')))


def _extract_adnc_grade(text) -> str:
    if pd.isna(text):
        return ''
    t = str(text).lower()
    if 'high' in t:
        return 'High'
    if 'intermediate' in t:
        return 'Intermediate'
    if 'low' in t:
        return 'Low'
    return ''


# ---------------------------------------------------------------------------
# Parse SMC xlsx
# ---------------------------------------------------------------------------

def _parse_smc_xlsx(path: Path) -> pd.DataFrame:
    df = pd.read_excel(path, engine='openpyxl', header=0)

    # Rename columns to stable keys
    col_map = {
        df.columns[0]: 'raw_id',
        df.columns[4]: 'age',
        df.columns[5]: 'sex',
        df.columns[6]: 'condition_raw',
        df.columns[8]: 'primary_dx',
        df.columns[9]: 'A_score',
        df.columns[10]: 'B_score',
        df.columns[11]: 'C_score',
    }
    df = df.rename(columns=col_map)

    # Keep SMC rows only (PO stubs in xlsx are almost entirely NaN — we use pptx instead)
    smc_mask = df['raw_id'].astype(str).str.match(r'^SMC-\d+$', na=False)
    df = df[smc_mask].copy()

    rows = []
    for _, r in df.iterrows():
        donor_id  = harmonize_id(r['raw_id'])
        condition = _norm_condition(r['condition_raw'])
        age       = float(r['age']) if not pd.isna(r['age']) else float('nan')
        sex_raw   = str(r['sex']).strip().upper() if not pd.isna(r['sex']) else ''
        sex       = sex_raw if sex_raw in ('M', 'F') else ''

        braak  = _b_to_braak(r['B_score'])
        thal   = _abc_to_int(r['A_score'])
        cerad  = _abc_to_int(r['C_score'])
        adnc   = _extract_adnc_grade(r['primary_dx'])
        apoe   = ''  # not collected for SMC cohort

        rows.append({
            'donor_id':   donor_id,
            'condition':  condition,
            'age':        age,
            'sex':        sex,
            'braak_stage': braak,
            'thal_phase': thal,
            'cerad_score': cerad,
            'apoe':       apoe,
            'adnc_grade': adnc,
        })

    return pd.DataFrame(rows)


# ---------------------------------------------------------------------------
# Parse PO pptx
# ---------------------------------------------------------------------------

def _parse_po_pptx(path: Path) -> pd.DataFrame:
    prs = Presentation(str(path))

    # Find the table shape (confirmed: 1 slide, 1 table)
    tbl = None
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                tbl = shape.table
                break
        if tbl:
            break

    if tbl is None:
        raise RuntimeError(f"No table found in {path}")

    # Extract raw cell text, skip 2 header rows
    raw_rows = []
    for i, row in enumerate(tbl.rows):
        if i < 2:
            continue
        raw_rows.append([cell.text_frame.text.strip() for cell in row.cells])

    # Columns after headers:
    # 0=자원종류, 1=Subject No., 2=병리진단명, 3=나이/성별, 4=APOE,
    # 5=ADNC description, 6=Thal phase, 7=Braak stage, 8=CERAD score

    # Handle P_011 split: merge consecutive rows where Subject No. is empty
    merged = []
    current = None
    for row in raw_rows:
        subj = row[1].strip()
        if subj:
            if current is not None:
                merged.append(current)
            current = row[:]
        else:
            # Fill empty fields in current from this continuation row
            if current is not None:
                for col in range(len(row)):
                    if not current[col] and row[col]:
                        current[col] = row[col]
    if current is not None:
        merged.append(current)

    rows = []
    for row in merged:
        raw_id  = row[1].strip()
        if not raw_id:
            continue
        donor_id  = harmonize_id(raw_id)
        cond_raw  = row[2].strip()
        condition = _norm_condition(cond_raw)

        # Parse "83/F" → age, sex
        age_sex = row[3].strip()
        age, sex = float('nan'), ''
        if '/' in age_sex:
            parts = age_sex.split('/')
            try:
                age = float(parts[0].strip())
            except ValueError:
                pass
            sex = parts[1].strip().upper() if len(parts) > 1 else ''

        apoe_raw = row[4].strip()
        apoe = '' if apoe_raw == '-' else apoe_raw

        adnc_text = row[5].strip()
        adnc      = _extract_adnc_grade(adnc_text)

        thal  = _roman_to_int(row[6])
        braak = _roman_to_int(row[7])
        cerad = _cerad_text_to_int(row[8])

        rows.append({
            'donor_id':    donor_id,
            'condition':   condition,
            'age':         age,
            'sex':         sex,
            'braak_stage': braak,
            'thal_phase':  thal,
            'cerad_score': cerad,
            'apoe':        apoe,
            'adnc_grade':  adnc,
        })

    return pd.DataFrame(rows)


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def run(qc_log: QCLogger) -> pd.DataFrame:
    smc_df = _parse_smc_xlsx(SMC_XLSX)
    po_df  = _parse_po_pptx(PO_PPTX)

    meta = pd.concat([smc_df, po_df], ignore_index=True)
    meta = meta.set_index('donor_id')

    # Log harmonization counts
    qc_log.log_id_harmonization(
        n_smc=len(smc_df), n_po=len(po_df), n_total=len(meta)
    )

    # Flag donors with missing critical fields
    critical = ['condition', 'age', 'sex', 'braak_stage']
    for donor_id, row in meta.iterrows():
        missing = [f for f in critical if pd.isna(row[f]) or row[f] == '']
        if missing:
            qc_log.log_missing_metadata(donor_id, missing)

    print(f"[Step 0.1] {len(smc_df)} SMC + {len(po_df)} PO donors = {len(meta)} total")
    return meta
